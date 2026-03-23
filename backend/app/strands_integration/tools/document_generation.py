"""
Document generation tools: DOCX, PPTX, XLSX, and chart (SVG) output.

All documents are uploaded to S3 (DOCUMENT_BUCKET) under agent-workspace/
and returned as a presigned URL valid for 1 hour.

Ported from aws-samples/sample-aws-idp-pipeline .skills/ definitions,
adapted to run as native Strands tools without code_interpreter/AgentCore.
"""

from __future__ import annotations

import io
import json
import logging
import math
import os
import re
import uuid
from datetime import datetime
from typing import Any

import boto3
from app.repositories.models.custom_bot import BotModel
from strands import tool
from strands.types.tools import AgentTool as StrandsAgentTool

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

DOCUMENT_BUCKET = os.environ.get("DOCUMENT_BUCKET", "")
REGION = os.environ.get("REGION", "us-east-1")
WORKSPACE_PREFIX = "agent-workspace"


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _s3_client():
    return boto3.client("s3", region_name=REGION)


def _upload_and_presign(
    data: bytes,
    filename: str,
    content_type: str,
    presign_expiry: int = 3600,
) -> str:
    """Upload bytes to S3 and return a 1-hour presigned download URL."""
    datestamp = datetime.utcnow().strftime("%Y/%m/%d")
    artifact_id = str(uuid.uuid4())
    key = f"{WORKSPACE_PREFIX}/documents/{datestamp}/{artifact_id}/{filename}"

    s3 = _s3_client()
    s3.put_object(
        Bucket=DOCUMENT_BUCKET,
        Key=key,
        Body=data,
        ContentType=content_type,
        ContentDisposition=f'attachment; filename="{filename}"',
    )
    logger.info(f"[DOC_GEN] Uploaded {filename} → s3://{DOCUMENT_BUCKET}/{key}")

    url = s3.generate_presigned_url(
        "get_object",
        Params={"Bucket": DOCUMENT_BUCKET, "Key": key},
        ExpiresIn=presign_expiry,
    )
    return url


def _parse_json_arg(value: str | list | dict, default: Any = None) -> Any:
    """Accept either a JSON string or a Python object; return parsed value."""
    if isinstance(value, (list, dict)):
        return value
    if isinstance(value, str):
        try:
            return json.loads(value)
        except (json.JSONDecodeError, ValueError):
            pass
    return default


# ---------------------------------------------------------------------------
# DOCX tool
# ---------------------------------------------------------------------------


def create_docx_tool(bot: BotModel | None = None) -> StrandsAgentTool:
    @tool
    def generate_docx(
        filename: str,
        content: str,
        title: str = "",
    ) -> str:
        """
        Generate a Word document (.docx) from structured Markdown-style content
        and return a presigned download URL (valid 1 hour).

        Supports the following Markdown conventions in `content`:
          - # Heading 1 / ## Heading 2 / ### Heading 3
          - **bold** and *italic* inline formatting
          - Unordered lists starting with "- " or "* "
          - Ordered lists starting with "1. ", "2. " etc.
          - Simple pipe tables: | Col A | Col B |\\n|---|---|\\n| val | val |
          - Blank lines between paragraphs

        Args:
            filename: Output filename, e.g. "report.docx". Must end with .docx.
            content:  Document body as Markdown-style text. The agent should
                      write well-structured content here.
            title:    Optional document title inserted as Heading 1 at the top.

        Returns:
            str: Presigned URL to download the generated document, or an error message.
        """
        if not DOCUMENT_BUCKET:
            return "Error: DOCUMENT_BUCKET environment variable is not set."

        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
        except ImportError:
            return (
                "Error: python-docx is not installed. "
                "Add 'python-docx' to pyproject.toml dependencies."
            )

        if not filename.lower().endswith(".docx"):
            filename = filename + ".docx"

        # --------------- ACIL Allen brand helpers ---------------
        # Colour palette (Purple hierarchy is primary)
        AA_PURPLE_2 = RGBColor(0x33, 0x10, 0x63)  # #331063 — sidebar / H1
        AA_PURPLE_3 = RGBColor(0x6C, 0x3F, 0x99)  # #6C3F99 — brand purple / H2–H3
        AA_GREY_2   = RGBColor(0x4D, 0x4D, 0x4D)  # #4D4D4D — body text
        AA_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

        def _style_heading(heading, level: int):
            """Apply YuGothic Medium + ACIL Allen purple to a heading paragraph."""
            color = AA_PURPLE_2 if level == 1 else AA_PURPLE_3
            for run in heading.runs:
                run.font.name = "Yu Gothic Medium"
                run.font.color.rgb = color

        def _set_cell_bg(cell, hex_color: str):
            """Set table cell background to a solid hex colour via OOXML."""
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), hex_color)
            tcPr.append(shd)
        # --------------------------------------------------------

        doc = Document()

        # Body font: Arial Narrow (brand-approved print font)
        doc.styles["Normal"].font.name = "Arial Narrow"

        # Optional title as H1
        if title:
            h = doc.add_heading(title, level=1)
            _style_heading(h, 1)

        def _add_run_with_inline(paragraph, text: str):
            """Parse **bold** and *italic* and add runs."""
            # Simple state-machine for ** and *
            pattern = re.compile(r"(\*\*.*?\*\*|\*.*?\*)")
            last = 0
            for m in pattern.finditer(text):
                if m.start() > last:
                    paragraph.add_run(text[last : m.start()])
                raw = m.group()
                if raw.startswith("**"):
                    run = paragraph.add_run(raw[2:-2])
                    run.bold = True
                else:
                    run = paragraph.add_run(raw[1:-1])
                    run.italic = True
                last = m.end()
            if last < len(text):
                paragraph.add_run(text[last:])

        def _is_table_line(line: str) -> bool:
            return line.strip().startswith("|") and line.strip().endswith("|")

        def _add_table(doc, table_lines: list[str]):
            rows = [
                [cell.strip() for cell in line.strip().strip("|").split("|")]
                for line in table_lines
                if not re.match(r"^\s*\|[-:| ]+\|\s*$", line)  # skip separator
            ]
            if not rows:
                return
            col_count = max(len(r) for r in rows)
            table = doc.add_table(rows=len(rows), cols=col_count)
            table.style = "Table Grid"
            for ri, row_data in enumerate(rows):
                for ci, cell_text in enumerate(row_data):
                    if ci < col_count:
                        cell = table.cell(ri, ci)
                        cell.text = cell_text
                        if ri == 0:
                            # ACIL Allen branded header: Purple 2 bg, white bold text
                            _set_cell_bg(cell, "331063")
                            for run in cell.paragraphs[0].runs:
                                run.bold = True
                                run.font.color.rgb = AA_WHITE
                                run.font.name = "Yu Gothic Medium"

        lines = content.splitlines()
        i = 0
        while i < len(lines):
            line = lines[i]

            # Headings
            if line.startswith("### "):
                h = doc.add_heading(line[4:].strip(), level=3)
                _style_heading(h, 3)
            elif line.startswith("## "):
                h = doc.add_heading(line[3:].strip(), level=2)
                _style_heading(h, 2)
            elif line.startswith("# "):
                h = doc.add_heading(line[2:].strip(), level=1)
                _style_heading(h, 1)

            # Table block
            elif _is_table_line(line):
                table_lines = []
                while i < len(lines) and _is_table_line(lines[i]):
                    table_lines.append(lines[i])
                    i += 1
                _add_table(doc, table_lines)
                continue

            # Unordered list
            elif re.match(r"^\s*[-*]\s+", line):
                para = doc.add_paragraph(style="List Bullet")
                _add_run_with_inline(para, re.sub(r"^\s*[-*]\s+", "", line))

            # Ordered list
            elif re.match(r"^\s*\d+\.\s+", line):
                para = doc.add_paragraph(style="List Number")
                _add_run_with_inline(para, re.sub(r"^\s*\d+\.\s+", "", line))

            # Blank line → skip
            elif line.strip() == "":
                pass

            # Normal paragraph
            else:
                para = doc.add_paragraph()
                _add_run_with_inline(para, line.strip())

            i += 1

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        url = _upload_and_presign(
            buf.read(),
            filename,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        return (
            f"Word document '{filename}' generated successfully.\n"
            f"Download (link valid 1 hour):\n{url}"
        )

    return generate_docx


# ---------------------------------------------------------------------------
# PPTX tool
# ---------------------------------------------------------------------------


def create_pptx_tool(bot: BotModel | None = None) -> StrandsAgentTool:
    @tool
    def generate_pptx(
        filename: str,
        title: str,
        slides: str,
    ) -> str:
        """
        Generate a PowerPoint presentation (.pptx) and return a presigned download URL.

        Args:
            filename: Output filename, e.g. "presentation.pptx". Must end with .pptx.
            title:    Presentation title shown on the title slide.
            slides:   JSON array of slide objects. Each object must have:
                        - "title" (str): Slide title
                        - "content" (list[str]): Bullet-point lines for the slide body.
                      Example:
                        [
                          {"title": "Overview", "content": ["Point A", "Point B"]},
                          {"title": "Details",  "content": ["Item 1", "Item 2", "Item 3"]}
                        ]

        Returns:
            str: Presigned URL to download the generated presentation, or an error message.
        """
        if not DOCUMENT_BUCKET:
            return "Error: DOCUMENT_BUCKET environment variable is not set."

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            return (
                "Error: python-pptx is not installed. "
                "Add 'python-pptx' to pyproject.toml dependencies."
            )

        slides_data = _parse_json_arg(slides, default=[])
        if not slides_data:
            return "Error: 'slides' must be a non-empty JSON array of slide objects."

        # --------------- ACIL Allen brand colours ---------------
        PPTX_PURPLE_1 = RGBColor(0x14, 0x00, 0x34)  # #140034 — title slide bg
        PPTX_PURPLE_2 = RGBColor(0x33, 0x10, 0x63)  # #331063 — content title
        PPTX_PURPLE_3 = RGBColor(0x6C, 0x3F, 0x99)  # #6C3F99 — accents
        PPTX_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
        PPTX_GREY_2   = RGBColor(0x4D, 0x4D, 0x4D)  # #4D4D4D — body text

        def _apply_font(text_frame, font_name: str, color: RGBColor | None = None):
            """Set font name (and optionally colour) on all paragraphs/runs."""
            for para in text_frame.paragraphs:
                para.font.name = font_name
                if color:
                    para.font.color.rgb = color
                for run in para.runs:
                    run.font.name = font_name
                    if color:
                        run.font.color.rgb = color
        # --------------------------------------------------------

        prs = Presentation()
        # Use widescreen 16:9
        prs.slide_width = int(10 * 914400)   # 10 inches
        prs.slide_height = int(5.625 * 914400)  # 5.625 inches

        slide_layouts = prs.slide_layouts

        # Title slide — dark Purple 1 background, white text
        title_slide_layout = slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        # Brand background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PPTX_PURPLE_1
        # Title text
        slide.shapes.title.text = title
        _apply_font(slide.shapes.title.text_frame, "Yu Gothic Medium", PPTX_WHITE)
        if slide.placeholders[1]:
            slide.placeholders[1].text = ""
            _apply_font(slide.placeholders[1].text_frame, "Arial Narrow", PPTX_WHITE)

        # Content slides
        content_layout = slide_layouts[1]  # Title and Content
        for slide_def in slides_data:
            s_title = slide_def.get("title", "")
            s_content = slide_def.get("content", [])

            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = s_title
            # Slide title: YuGothic Medium, Purple 2
            _apply_font(s.shapes.title.text_frame, "Yu Gothic Medium", PPTX_PURPLE_2)

            tf = s.placeholders[1].text_frame
            tf.clear()
            for idx, bullet in enumerate(s_content):
                if idx == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            # Body text: Arial Narrow, Grey 2
            _apply_font(tf, "Arial Narrow", PPTX_GREY_2)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        if not filename.lower().endswith(".pptx"):
            filename = filename + ".pptx"

        url = _upload_and_presign(
            buf.read(),
            filename,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        return (
            f"PowerPoint presentation '{filename}' generated successfully "
            f"({len(slides_data)} content slides + title slide).\n"
            f"Download (link valid 1 hour):\n{url}"
        )

    return generate_pptx


# ---------------------------------------------------------------------------
# XLSX tool
# ---------------------------------------------------------------------------


def create_xlsx_tool(bot: BotModel | None = None) -> StrandsAgentTool:
    @tool
    def generate_xlsx(
        filename: str,
        sheets: str,
    ) -> str:
        """
        Generate an Excel spreadsheet (.xlsx) and return a presigned download URL.

        Args:
            filename: Output filename, e.g. "data.xlsx". Must end with .xlsx.
            sheets:   JSON array of sheet objects. Each object must have:
                        - "name" (str): Sheet tab name
                        - "headers" (list[str]): Column header labels
                        - "rows" (list[list]): Data rows, each a list of cell values
                      Optionally:
                        - "title" (str): Bold title row inserted above headers
                      Example:
                        [
                          {
                            "name": "Sales",
                            "title": "Q1 Sales Report",
                            "headers": ["Region", "Product", "Revenue"],
                            "rows": [
                              ["North", "Widget A", 12500],
                              ["South", "Widget B", 9800]
                            ]
                          }
                        ]

        Returns:
            str: Presigned URL to download the generated spreadsheet, or an error message.
        """
        if not DOCUMENT_BUCKET:
            return "Error: DOCUMENT_BUCKET environment variable is not set."

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        except ImportError:
            return (
                "Error: openpyxl is not installed. "
                "Add 'openpyxl' to pyproject.toml dependencies."
            )

        sheets_data = _parse_json_arg(sheets, default=[])
        if not sheets_data:
            return "Error: 'sheets' must be a non-empty JSON array of sheet objects."

        wb = Workbook()
        wb.remove(wb.active)  # remove default empty sheet

        # ACIL Allen brand colours — Purple 2 (#331063) headers, Arial Narrow body
        HEADER_FILL  = PatternFill("solid", fgColor="331063")   # AA Purple 2
        TITLE_FONT   = Font(bold=True, size=13, name="Yu Gothic Medium", color="331063")
        HEADER_FONT  = Font(bold=True, color="FFFFFF", name="Yu Gothic Medium")
        BODY_FONT    = Font(name="Arial Narrow")
        HEADER_ALIGN = Alignment(horizontal="center", vertical="center")
        THIN = Side(style="thin", color="C8C8C8")               # AA Grey 4
        CELL_BORDER  = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

        for sheet_def in sheets_data:
            sheet_name = sheet_def.get("name", "Sheet")[:31]  # Excel max 31 chars
            headers = sheet_def.get("headers", [])
            rows = sheet_def.get("rows", [])
            sheet_title = sheet_def.get("title", "")

            ws = wb.create_sheet(title=sheet_name)
            current_row = 1

            # Optional bold title row
            if sheet_title:
                ws.cell(row=current_row, column=1, value=sheet_title).font = TITLE_FONT
                current_row += 1

            # Header row
            if headers:
                for col_idx, header in enumerate(headers, start=1):
                    cell = ws.cell(row=current_row, column=col_idx, value=header)
                    cell.font = HEADER_FONT
                    cell.fill = HEADER_FILL
                    cell.alignment = HEADER_ALIGN
                    cell.border = CELL_BORDER
                current_row += 1

            # Data rows
            for row_data in rows:
                for col_idx, value in enumerate(row_data, start=1):
                    cell = ws.cell(row=current_row, column=col_idx, value=value)
                    cell.border = CELL_BORDER
                    cell.font = BODY_FONT
                current_row += 1

            # Auto-fit column widths (approximate)
            for col_cells in ws.columns:
                max_len = 0
                col_letter = col_cells[0].column_letter
                for c in col_cells:
                    if c.value is not None:
                        max_len = max(max_len, len(str(c.value)))
                ws.column_dimensions[col_letter].width = min(max_len + 4, 60)

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)

        if not filename.lower().endswith(".xlsx"):
            filename = filename + ".xlsx"

        url = _upload_and_presign(
            buf.read(),
            filename,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        return (
            f"Excel spreadsheet '{filename}' generated successfully "
            f"({len(sheets_data)} sheet(s)).\n"
            f"Download (link valid 1 hour):\n{url}"
        )

    return generate_xlsx


# ---------------------------------------------------------------------------
# Chart (SVG) tool — pure Python, no binary dependencies
# ---------------------------------------------------------------------------

# ACIL Allen brand colour palette — hierarchy: Purple → Grey → Blue → Green
_DEFAULT_COLORS = [
    "#6C3F99",  # AA Purple 3 (primary brand)
    "#006A9F",  # AA Blue 3
    "#428D52",  # AA Green 3
    "#7BBDD6",  # AA Blue 4
    "#331063",  # AA Purple 2
    "#9FD18B",  # AA Green 4
    "#9D85BE",  # AA Purple 4
    "#4D4D4D",  # AA Grey 2
]

_W, _H = 800, 480
_PAD_TOP, _PAD_BOTTOM, _PAD_LEFT, _PAD_RIGHT = 55, 75, 75, 30


def _xe(s: str) -> str:
    """Escape XML special characters."""
    return str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


def _nice_ticks(lo: float, hi: float, n: int = 5) -> list[float]:
    """Return ~n human-friendly tick values covering [lo, hi]."""
    if hi == lo:
        hi = lo + 1
    raw = (hi - lo) / n
    mag = 10 ** math.floor(math.log10(raw))
    step = mag * min((1, 2, 5, 10), key=lambda m: abs(m * mag - raw))
    start = math.floor(lo / step) * step
    ticks = []
    v = start
    while v <= hi + step * 0.01:
        ticks.append(round(v, 10))
        v += step
    return ticks


def _svg_chart(
    chart_type: str,
    title: str,
    labels_data: list,
    values_data: list,
    x_label: str,
    y_label: str,
    colors_data: list,
) -> bytes:
    W, H = _W, _H
    PT, PB, PL, PR = _PAD_TOP, _PAD_BOTTOM, _PAD_LEFT, _PAD_RIGHT
    CW, CH = W - PL - PR, H - PT - PB
    els: list[str] = []

    def g(tag: str, **attrs) -> str:
        parts = " ".join(f'{k.replace("_", "-")}="{v}"' for k, v in attrs.items())
        return f"<{tag} {parts}/>"

    def txt(x, y, s, **kw) -> str:
        # ACIL Allen brand font: Arial Narrow; default text colour: AA Grey 2
        kw.setdefault("font_family", "'Arial Narrow',Arial,sans-serif")
        kw.setdefault("fill", "#4D4D4D")  # AA Grey 2
        parts = " ".join(f'{k.replace("_", "-")}="{v}"' for k, v in kw.items())
        return f'<text x="{x}" y="{y}" {parts}>{_xe(s)}</text>'

    # White background
    els.append(g("rect", width=W, height=H, fill="#FFFFFF"))
    # Title — YuGothic Medium approximated via font-weight bold; AA Purple 2 colour
    els.append(txt(W // 2, 38, title, text_anchor="middle", font_size=18,
                   font_weight="bold", fill="#331063"))

    if chart_type in ("pie", "donut"):
        total = sum(float(v) for v in values_data)
        if total == 0:
            total = 1
        cx, cy = W // 2, PT + CH // 2
        r = min(CW, CH) // 2 - 10
        inner_r = r // 2 if chart_type == "donut" else 0
        angle = -90.0
        for i, (lbl, val) in enumerate(zip(labels_data, values_data)):
            sweep = float(val) / total * 360
            a1 = math.radians(angle)
            a2 = math.radians(angle + sweep)
            x1, y1 = cx + r * math.cos(a1), cy + r * math.sin(a1)
            x2, y2 = cx + r * math.cos(a2), cy + r * math.sin(a2)
            large = 1 if sweep > 180 else 0
            color = colors_data[i % len(colors_data)]
            if inner_r > 0:
                ix1, iy1 = cx + inner_r * math.cos(a1), cy + inner_r * math.sin(a1)
                ix2, iy2 = cx + inner_r * math.cos(a2), cy + inner_r * math.sin(a2)
                d = (f"M {x1:.2f} {y1:.2f} A {r} {r} 0 {large} 1 {x2:.2f} {y2:.2f} "
                     f"L {ix2:.2f} {iy2:.2f} A {inner_r} {inner_r} 0 {large} 0 {ix1:.2f} {iy1:.2f} Z")
            else:
                d = (f"M {cx} {cy} L {x1:.2f} {y1:.2f} "
                     f"A {r} {r} 0 {large} 1 {x2:.2f} {y2:.2f} Z")
            els.append(f'<path d="{d}" fill="{color}" stroke="#fff" stroke-width="1.5"/>')
            # Label at midpoint of arc
            mid_a = math.radians(angle + sweep / 2)
            lx = cx + (r * 0.7 + inner_r * 0.3) * math.cos(mid_a)
            ly = cy + (r * 0.7 + inner_r * 0.3) * math.sin(mid_a)
            pct = f"{float(val)/total*100:.1f}%"
            els.append(txt(f"{lx:.1f}", f"{ly:.1f}", pct, text_anchor="middle",
                           font_size=11, fill="#fff", font_weight="bold"))
            angle += sweep
        # Legend (right side)
        lx0 = cx + r + 20
        for i, lbl in enumerate(labels_data):
            ly0 = PT + 20 + i * 22
            color = colors_data[i % len(colors_data)]
            els.append(g("rect", x=lx0, y=ly0, width=14, height=14,
                         fill=color, rx=2))
            els.append(txt(lx0 + 18, ly0 + 11, lbl, font_size=12))

    else:
        # Axis lines — AA Grey 3 (#7D7D7D)
        els.append(g("line", x1=PL, y1=PT, x2=PL, y2=H - PB,
                     stroke="#7D7D7D", stroke_width="1.5"))
        els.append(g("line", x1=PL, y1=H - PB, x2=W - PR, y2=H - PB,
                     stroke="#7D7D7D", stroke_width="1.5"))

        # Axis labels — AA Grey 2 (#4D4D4D)
        if x_label:
            els.append(txt(PL + CW // 2, H - 8, x_label,
                           text_anchor="middle", font_size=12, fill="#4D4D4D"))
        if y_label:
            mid = PT + CH // 2
            els.append(
                f'<text x="14" y="{mid}" text-anchor="middle" '
                f'font-family="\'Arial Narrow\',Arial,sans-serif" font-size="12" fill="#4D4D4D" '
                f'transform="rotate(-90,14,{mid})">{_xe(y_label)}</text>'
            )

        is_multi = bool(values_data) and isinstance(values_data[0], dict)

        if chart_type == "horizontal_bar":
            all_vals = [float(v) for v in values_data]
            max_v = max(all_vals) if all_vals else 1
            ticks = _nice_ticks(0, max_v)
            tick_max = ticks[-1]
            n = len(labels_data)
            bar_h = CH / max(n, 1) * 0.6
            gap = CH / max(n, 1)

            for tick in ticks:
                tx = PL + CW * tick / tick_max
                els.append(g("line", x1=tx, y1=PT, x2=tx, y2=H - PB,
                             stroke="#C8C8C8", stroke_width="1"))
                els.append(txt(f"{tx:.1f}", H - PB + 16, f"{tick:g}",
                               text_anchor="middle", font_size=10, fill="#7D7D7D"))

            for i, (lbl, val) in enumerate(zip(labels_data, all_vals)):
                by = PT + i * gap + (gap - bar_h) / 2
                bw = CW * val / tick_max
                color = colors_data[i % len(colors_data)]
                els.append(g("rect", x=PL, y=f"{by:.1f}", width=f"{bw:.1f}",
                             height=f"{bar_h:.1f}", fill=color, rx=2, opacity="0.85"))
                els.append(txt(PL - 6, f"{by + bar_h/2 + 4:.1f}", lbl,
                               text_anchor="end", font_size=11, fill="#4D4D4D"))

        elif chart_type == "bar":
            if is_multi:
                series_list = values_data
                all_flat = [float(x) for s in series_list for x in s["data"]]
            else:
                series_list = [{"label": "", "data": values_data}]
                all_flat = [float(v) for v in values_data]

            max_v = max(all_flat) if all_flat else 1
            ticks = _nice_ticks(0, max_v)
            tick_max = ticks[-1]

            for tick in ticks:
                ty = H - PB - CH * tick / tick_max
                els.append(g("line", x1=PL, y1=f"{ty:.1f}", x2=W - PR, y2=f"{ty:.1f}",
                             stroke="#C8C8C8", stroke_width="1"))
                els.append(txt(PL - 6, f"{ty + 4:.1f}", f"{tick:g}",
                               text_anchor="end", font_size=10, fill="#7D7D7D"))

            n_groups = len(labels_data)
            n_series = len(series_list)
            group_w = CW / max(n_groups, 1)
            bar_w = group_w * 0.8 / n_series

            for gi, lbl in enumerate(labels_data):
                gx = PL + gi * group_w + group_w * 0.1
                for si, series in enumerate(series_list):
                    val = float(series["data"][gi]) if gi < len(series["data"]) else 0
                    bh = CH * val / tick_max
                    bx = gx + si * bar_w
                    by = H - PB - bh
                    color = colors_data[si % len(colors_data)]
                    els.append(g("rect", x=f"{bx:.1f}", y=f"{by:.1f}",
                                 width=f"{bar_w:.1f}", height=f"{bh:.1f}",
                                 fill=color, rx=2, opacity="0.85"))
                cx_lbl = PL + gi * group_w + group_w / 2
                els.append(txt(f"{cx_lbl:.1f}", H - PB + 16, lbl,
                               text_anchor="middle", font_size=10, fill="#555"))

            if is_multi:
                for si, series in enumerate(series_list):
                    color = colors_data[si % len(colors_data)]
                    lx0 = PL + si * 120
                    els.append(g("rect", x=lx0, y=H - 18, width=12, height=12,
                                 fill=color, rx=2))
                    els.append(txt(lx0 + 16, H - 7, series.get("label", f"Series {si+1}"),
                                   font_size=11, fill="#4D4D4D"))

        elif chart_type == "line":
            if is_multi:
                series_list = values_data
                all_flat = [float(x) for s in series_list for x in s["data"]]
            else:
                series_list = [{"label": "", "data": values_data}]
                all_flat = [float(v) for v in values_data]

            max_v = max(all_flat) if all_flat else 1
            ticks = _nice_ticks(0, max_v)
            tick_max = ticks[-1]
            n = len(labels_data)

            for tick in ticks:
                ty = H - PB - CH * tick / tick_max
                els.append(g("line", x1=PL, y1=f"{ty:.1f}", x2=W - PR, y2=f"{ty:.1f}",
                             stroke="#C8C8C8", stroke_width="1"))
                els.append(txt(PL - 6, f"{ty + 4:.1f}", f"{tick:g}",
                               text_anchor="end", font_size=10, fill="#7D7D7D"))

            for gi, lbl in enumerate(labels_data):
                px = PL + gi * CW / max(n - 1, 1)
                els.append(txt(f"{px:.1f}", H - PB + 16, lbl,
                               text_anchor="middle", font_size=10, fill="#555"))

            for si, series in enumerate(series_list):
                color = colors_data[si % len(colors_data)]
                pts = []
                for gi, val in enumerate(series["data"]):
                    px = PL + gi * CW / max(n - 1, 1)
                    py = H - PB - CH * float(val) / tick_max
                    pts.append((px, py))
                if pts:
                    pts_str = " ".join(f"{x:.1f},{y:.1f}" for x, y in pts)
                    els.append(f'<polyline points="{pts_str}" fill="none" '
                               f'stroke="{color}" stroke-width="2.5" stroke-linejoin="round"/>')
                    for px, py in pts:
                        els.append(g("circle", cx=f"{px:.1f}", cy=f"{py:.1f}", r=4,
                                     fill=color, stroke="#fff", stroke_width="1.5"))

            if is_multi:
                for si, series in enumerate(series_list):
                    color = colors_data[si % len(colors_data)]
                    lx0 = PL + si * 140
                    els.append(g("line", x1=lx0, y1=H - 12, x2=lx0 + 20, y2=H - 12,
                                 stroke=color, stroke_width="2.5"))
                    els.append(txt(lx0 + 24, H - 8, series.get("label", f"Series {si+1}"),
                                   font_size=11, fill="#4D4D4D"))

        elif chart_type == "scatter":
            xs = [float(v) for v in labels_data]
            ys = [float(v) for v in values_data]
            x_ticks = _nice_ticks(min(xs), max(xs))
            y_ticks = _nice_ticks(min(ys), max(ys))
            x_min, x_max = x_ticks[0], x_ticks[-1]
            y_min, y_max = y_ticks[0], y_ticks[-1]

            for tick in y_ticks:
                ty = H - PB - CH * (tick - y_min) / (y_max - y_min)
                els.append(g("line", x1=PL, y1=f"{ty:.1f}", x2=W - PR, y2=f"{ty:.1f}",
                             stroke="#C8C8C8", stroke_width="1"))
                els.append(txt(PL - 6, f"{ty + 4:.1f}", f"{tick:g}",
                               text_anchor="end", font_size=10, fill="#7D7D7D"))

            for tick in x_ticks:
                tx = PL + CW * (tick - x_min) / (x_max - x_min)
                els.append(g("line", x1=f"{tx:.1f}", y1=PT, x2=f"{tx:.1f}", y2=H - PB,
                             stroke="#C8C8C8", stroke_width="1"))
                els.append(txt(f"{tx:.1f}", H - PB + 16, f"{tick:g}",
                               text_anchor="middle", font_size=10, fill="#7D7D7D"))

            color = colors_data[0]
            for x, y in zip(xs, ys):
                px = PL + CW * (x - x_min) / (x_max - x_min)
                py = H - PB - CH * (y - y_min) / (y_max - y_min)
                els.append(g("circle", cx=f"{px:.1f}", cy=f"{py:.1f}", r=5,
                             fill=color, opacity="0.75", stroke="#fff", stroke_width="1"))

    svg = (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" '
        f'viewBox="0 0 {W} {H}">'
        + "".join(els)
        + "</svg>"
    )
    return svg.encode("utf-8")


def create_chart_tool(bot: BotModel | None = None) -> StrandsAgentTool:
    @tool
    def generate_chart(
        filename: str,
        chart_type: str,
        title: str,
        labels: str,
        values: str,
        x_label: str = "",
        y_label: str = "",
        colors: str = "",
    ) -> str:
        """
        Generate a data chart as an SVG image and return a presigned download URL.

        Supported chart types: bar, horizontal_bar, line, pie, donut, scatter.

        Args:
            filename:    Output filename, e.g. "sales_chart.svg". Must end with .svg.
            chart_type:  One of: bar, horizontal_bar, line, pie, donut, scatter.
            title:       Chart title.
            labels:      JSON array of category labels, e.g. ["Jan","Feb","Mar"].
                         For scatter charts: JSON array of numeric x-values.
            values:      JSON array of numeric values, e.g. [120, 95, 140].
                         For multi-series bar/line: JSON array of series objects,
                         e.g. [{"label":"Series A","data":[1,2,3]},{"label":"Series B","data":[4,5,6]}]
            x_label:     Optional X-axis label (bar, line, scatter only).
            y_label:     Optional Y-axis label (bar, line, scatter only).
            colors:      Optional JSON array of hex color strings, e.g. ["#2196F3","#FF5722"].
                         Falls back to a professional default palette if not provided.

        Returns:
            str: Presigned URL to view/download the chart SVG, or an error message.
        """
        if not DOCUMENT_BUCKET:
            return "Error: DOCUMENT_BUCKET environment variable is not set."

        VALID_TYPES = {"bar", "horizontal_bar", "line", "pie", "donut", "scatter"}
        chart_type = chart_type.lower().strip()
        if chart_type not in VALID_TYPES:
            return f"Error: chart_type must be one of {sorted(VALID_TYPES)}."

        labels_data = _parse_json_arg(labels, default=[])
        values_data = _parse_json_arg(values, default=[])
        colors_data = _parse_json_arg(colors, default=[]) or _DEFAULT_COLORS

        if not labels_data:
            return "Error: 'labels' must be a non-empty JSON array."
        if not values_data:
            return "Error: 'values' must be a non-empty JSON array."

        if not filename.lower().endswith(".svg"):
            filename = re.sub(r"\.[^.]+$", "", filename) + ".svg"

        svg_bytes = _svg_chart(
            chart_type, title, labels_data, values_data, x_label, y_label, colors_data
        )

        url = _upload_and_presign(svg_bytes, filename, "image/svg+xml")
        return (
            f"Chart '{filename}' ({chart_type}) generated successfully.\n"
            f"View/download (link valid 1 hour):\n{url}"
        )

    return generate_chart


# ---------------------------------------------------------------------------
# Suite factory
# ---------------------------------------------------------------------------


def create_document_generation_tools(bot: BotModel | None = None) -> list[StrandsAgentTool]:
    """Return all document-generation tools as a list."""
    return [
        create_docx_tool(bot),
        create_pptx_tool(bot),
        create_xlsx_tool(bot),
        create_chart_tool(bot),
    ]
